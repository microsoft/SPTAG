#!/usr/bin/env python3

from __future__ import annotations

from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TITLE_FONT = "Aptos"
BODY_FONT = "Aptos"
MONO_FONT = "Cascadia Code"


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip().lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


WHITE = rgb("FFFFFF")
BLACK = rgb("111111")
GRAY = rgb("6B7280")
LINE = rgb("222222")
PANEL = rgb("F3F4F6")
HEADER_FILL = rgb("111111")
HEADER_TEXT = rgb("FFFFFF")


def set_background(slide, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=WHITE, line_color=LINE, line_width=1.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines,
    font_name=BODY_FONT,
    font_size=18,
    color=BLACK,
    bold=False,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
    margins=(0.08, 0.05, 0.08, 0.05),
    line_spacing=1.1,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = vertical
    text_frame.margin_left = Inches(margins[0])
    text_frame.margin_top = Inches(margins[1])
    text_frame.margin_right = Inches(margins[2])
    text_frame.margin_bottom = Inches(margins[3])

    prepared_lines = [lines] if isinstance(lines, str) else list(lines)
    for index, line in enumerate(prepared_lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(4)
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_header(slide, title: str, subtitle: str = "") -> None:
    set_background(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.58), HEADER_FILL, HEADER_FILL, 0)
    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.13),
        Inches(10.6),
        Inches(0.3),
        title,
        font_name=TITLE_FONT,
        font_size=24,
        color=HEADER_TEXT,
        bold=True,
        margins=(0, 0, 0, 0),
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.78),
            Inches(0.82),
            Inches(11.0),
            Inches(0.3),
            subtitle,
            font_name=BODY_FONT,
            font_size=12,
            color=GRAY,
            margins=(0, 0, 0, 0),
        )
    add_shape(slide, Inches(0.78), Inches(1.2), Inches(11.8), Inches(0.02), LINE, LINE, 0)


def add_footer(slide, page_number: int, total_pages: int = 4) -> None:
    add_textbox(
        slide,
        Inches(11.7),
        Inches(7.0),
        Inches(0.5),
        Inches(0.18),
        f"{page_number:02d}/{total_pages:02d}",
        font_name=BODY_FONT,
        font_size=10,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.RIGHT,
        margins=(0, 0, 0, 0),
    )


def add_panel(slide, left, top, width, height, title: str, lines, fill_color=PANEL) -> None:
    add_shape(slide, left, top, width, height, fill_color, LINE, 1.0)
    add_textbox(
        slide,
        left + Inches(0.14),
        top + Inches(0.12),
        width - Inches(0.28),
        Inches(0.28),
        title,
        font_name=TITLE_FONT,
        font_size=16,
        color=BLACK,
        bold=True,
        margins=(0, 0, 0, 0),
    )
    add_textbox(
        slide,
        left + Inches(0.14),
        top + Inches(0.46),
        width - Inches(0.28),
        height - Inches(0.56),
        lines,
        font_name=BODY_FONT,
        font_size=14,
        color=BLACK,
        margins=(0, 0, 0, 0),
        line_spacing=1.15,
    )


def add_diagram_node(slide, left, top, width, height, text: str, fill_color=WHITE) -> None:
    add_shape(slide, left, top, width, height, fill_color, LINE, 1.2)
    add_textbox(
        slide,
        left + Inches(0.05),
        top + Inches(0.03),
        width - Inches(0.1),
        height - Inches(0.06),
        text,
        font_name=BODY_FONT,
        font_size=13,
        color=BLACK,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margins=(0, 0, 0, 0),
    )


def add_vertical_line(slide, left, top, height, color=LINE) -> None:
    add_shape(slide, left, top, Inches(0.02), height, color, color, 0)


def add_horizontal_line(slide, left, top, width, color=LINE) -> None:
    add_shape(slide, left, top, width, Inches(0.02), color, color, 0)


def add_bullets(slide, left, top, width, height, lines, font_size=18) -> None:
    add_textbox(
        slide,
        left,
        top,
        width,
        height,
        [f"- {line}" for line in lines],
        font_name=BODY_FONT,
        font_size=font_size,
        color=BLACK,
        margins=(0, 0, 0, 0),
        line_spacing=1.18,
    )


def add_simple_table(slide, left, top, col_widths, row_heights, values) -> None:
    current_top = top
    for row_index, row in enumerate(values):
        current_left = left
        for col_index, cell_value in enumerate(row):
            fill = PANEL if row_index == 0 else WHITE
            add_shape(slide, current_left, current_top, col_widths[col_index], row_heights[row_index], fill, LINE, 1.0)
            add_textbox(
                slide,
                current_left + Inches(0.04),
                current_top + Inches(0.03),
                col_widths[col_index] - Inches(0.08),
                row_heights[row_index] - Inches(0.06),
                cell_value,
                font_name=BODY_FONT,
                font_size=12 if row_index == 0 else 13,
                color=BLACK,
                bold=(row_index == 0 or col_index == 0),
                align=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE,
                margins=(0, 0, 0, 0),
            )
            current_left += col_widths[col_index]
        current_top += row_heights[row_index]


def build_slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Tag Tree Structure", "Current tag levels in the benchmark and serving logic")

    add_panel(
        slide,
        Inches(0.78),
        Inches(1.55),
        Inches(3.3),
        Inches(4.8),
        "Key points",
        [
            "Current levels: org / dept / team / project.",
            "The hierarchy goes from coarse to fine.",
            "Query filters may appear at any level.",
        ],
    )

    root_left = Inches(6.0)
    root_top = Inches(1.8)
    node_w = Inches(1.4)
    node_h = Inches(0.48)

    add_diagram_node(slide, root_left, root_top, node_w, node_h, "org_*")
    add_vertical_line(slide, root_left + Inches(0.69), root_top + node_h, Inches(0.25))
    add_horizontal_line(slide, root_left - Inches(1.3), root_top + Inches(0.73), Inches(4.0))

    dept_y = Inches(2.55)
    dept_xs = [Inches(4.85), Inches(7.15)]
    for x in dept_xs:
        add_vertical_line(slide, x + Inches(0.69), Inches(2.53), Inches(0.22))
        add_diagram_node(slide, x, dept_y, node_w, node_h, "dept_*")

    add_vertical_line(slide, Inches(5.54), dept_y + node_h, Inches(0.22))
    add_vertical_line(slide, Inches(7.84), dept_y + node_h, Inches(0.22))
    add_horizontal_line(slide, Inches(4.05), Inches(3.25), Inches(5.3))

    team_y = Inches(3.45)
    team_xs = [Inches(3.35), Inches(5.15), Inches(6.95), Inches(8.75)]
    for x in team_xs:
        add_vertical_line(slide, x + Inches(0.69), Inches(3.23), Inches(0.22))
        add_diagram_node(slide, x, team_y, node_w, node_h, "team_*")

    add_vertical_line(slide, Inches(4.04), team_y + node_h, Inches(0.22))
    add_vertical_line(slide, Inches(9.44), team_y + node_h, Inches(0.22))
    add_horizontal_line(slide, Inches(4.04), Inches(4.15), Inches(5.4))

    project_y = Inches(4.35)
    project_xs = [Inches(3.35), Inches(8.75)]
    for x in project_xs:
        add_vertical_line(slide, x + Inches(0.69), Inches(4.13), Inches(0.22))
        add_diagram_node(slide, x, project_y, node_w, node_h, "project_*")

    add_textbox(
        slide,
        Inches(4.5),
        Inches(5.55),
        Inches(4.3),
        Inches(0.25),
        "coarse -> fine",
        font_name=BODY_FONT,
        font_size=12,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )
    add_footer(slide, 1)


def build_slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Sparse / Dense Strategy Selection", "Route choice uses exact posting fanout, not tenant-size heuristics")

    add_diagram_node(slide, Inches(4.8), Inches(1.65), Inches(3.2), Inches(0.5), "query + tenant + tags", PANEL)
    add_vertical_line(slide, Inches(6.39), Inches(2.15), Inches(0.24))
    add_diagram_node(slide, Inches(3.5), Inches(2.45), Inches(5.8), Inches(0.8), "all tags have direct posting lists and exact union <= DirectSparseMaxPostings ?")
    add_vertical_line(slide, Inches(6.39), Inches(3.25), Inches(0.2))
    add_horizontal_line(slide, Inches(2.35), Inches(3.75), Inches(8.1))
    add_vertical_line(slide, Inches(3.1), Inches(3.75), Inches(0.25))
    add_vertical_line(slide, Inches(9.85), Inches(3.75), Inches(0.25))
    add_textbox(slide, Inches(2.5), Inches(3.5), Inches(0.8), Inches(0.2), "yes", font_name=BODY_FONT, font_size=11, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(9.5), Inches(3.5), Inches(1.0), Inches(0.2), "no", font_name=BODY_FONT, font_size=11, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))

    add_panel(
        slide,
        Inches(0.78),
        Inches(4.1),
        Inches(5.55),
        Inches(2.05),
        "Sparse path algorithm",
        [
            "Collect the exact direct posting union for all query tags.",
            "Bypass graph search and scan those postings directly.",
            "Keep exact inline tag filtering during posting scan.",
        ],
    )

    add_panel(
        slide,
        Inches(6.98),
        Inches(4.1),
        Inches(5.55),
        Inches(2.05),
        "Dense path algorithm",
        [
            "Estimate selectivity from routing stats and head metadata.",
            "Run graph search and compute adaptive postingTarget.",
            "Keep exact inline tag filtering on the result path.",
        ],
    )

    add_textbox(
        slide,
        Inches(0.85),
        Inches(6.45),
        Inches(11.5),
        Inches(0.22),
        "Override: ForceDenseTagSearch sends all filtered queries to the dense path.",
        font_name=BODY_FONT,
        font_size=12,
        color=GRAY,
        bold=True,
        margins=(0, 0, 0, 0),
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 2)


def build_slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Adaptive Nprobe Algorithm", "Corrected to match the current wrapper + SearchIndex code path")

    add_panel(
        slide,
        Inches(0.78),
        Inches(1.62),
        Inches(3.35),
        Inches(1.3),
        "Wrapper step",
        [
            "Estimate raw filtered vector selectivity.",
            "sel_eff = clamp(sel_raw / max(1, safety), 1e-6, 1.0)",
        ],
    )

    add_panel(
        slide,
        Inches(4.45),
        Inches(1.62),
        Inches(8.05),
        Inches(3.65),
        "SearchIndex step",
        [
            "n_base = max(SearchInternalResultNum, topk)",
            "avgPosting = max(1, tenantSize / postingCount)",
            "p_recall = ceil(topk * R / max(1e-6, avgPosting * sel_eff))",
            "p_cover = ceil(n_base / max(1e-6, sel_eff^gamma))",
            "postingTarget = min(numSamples, max(n_base, p_recall, p_cover))",
            "graphResultNum = postingTarget",
        ],
        fill_color=WHITE,
    )

    add_panel(
        slide,
        Inches(0.78),
        Inches(3.2),
        Inches(3.35),
        Inches(2.05),
        "Interpretation",
        [
            "Lower selectivity means more postings are needed.",
            "The budget expands for both recall and coverage.",
            "Dense search uses postingTarget as graphResultNum.",
        ],
    )

    add_textbox(
        slide,
        Inches(4.7),
        Inches(5.52),
        Inches(7.55),
        Inches(0.22),
        "R = clamp(targetRecall, 0.01, 1.0)   |   gamma = clamp(coverageExponent, 0.0, 2.0)",
        font_name=MONO_FONT,
        font_size=10,
        color=GRAY,
        bold=True,
        margins=(0, 0, 0, 0),
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3)


def build_slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Sequence / Random Tag Filtering Performance", "Validated cache-stress snapshot")

    table_values = [
        ["Scenario", "Avg Latency", "P95", "P99", "QPS", "Recall"],
        ["sequence", "31.87 ms", "74.85 ms", "84.85 ms", "31.38", "0.9953"],
        ["random", "47.27 ms", "122.46 ms", "288.86 ms", "21.15", "0.9945"],
        ["delta", "+48.32%", "+63.61%", "+240.44%", "-32.60%", "-0.0008"],
    ]
    col_widths = [Inches(2.0), Inches(1.9), Inches(1.6), Inches(1.6), Inches(1.4), Inches(1.4)]
    row_heights = [Inches(0.55), Inches(0.65), Inches(0.65), Inches(0.65)]
    add_simple_table(slide, Inches(0.95), Inches(1.85), col_widths, row_heights, table_values)

    add_panel(
        slide,
        Inches(0.95),
        Inches(5.2),
        Inches(11.0),
        Inches(1.1),
        "Readout",
        [
            "Recall stays effectively stable.",
            "Random tenant/tag mixing mainly hurts tail latency and throughput.",
        ],
    )
    add_footer(slide, 4)


def build_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    build_slide_1(presentation)
    build_slide_2(presentation)
    build_slide_3(presentation)
    build_slide_4(presentation)
    return presentation


def main() -> None:
    pptx_path = Path(__file__).with_name("cache-acl-adaptive-nprobe-overview.pptx")
    backup_path = pptx_path.with_suffix(pptx_path.suffix + ".prettify.bak")
    if pptx_path.exists() and not backup_path.exists():
        shutil.copy2(pptx_path, backup_path)

    presentation = build_presentation()
    presentation.save(str(pptx_path))

    print(f"Updated {pptx_path}")
    print(f"Backup  {backup_path}")


if __name__ == "__main__":
    main()