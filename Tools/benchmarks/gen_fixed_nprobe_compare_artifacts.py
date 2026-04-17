#!/usr/bin/env python3

import argparse
import json
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


LEVEL_ORDER = ["org", "dept", "team", "project"]
LEVEL_LABELS = {
    "org": "Org",
    "dept": "Dept",
    "team": "Team",
    "project": "Project",
}
VARIANT_ORDER = ["old_single", "v5_headbundle"]
VARIANT_LABELS = {
    "old_single": "Original single-index",
    "v5_headbundle": "Current head-bundle",
}

TITLE_BG = RGBColor(0x10, 0x22, 0x3A)
ACCENT_BLUE = RGBColor(0x1F, 0x5A, 0xFF)
ACCENT_ORANGE = RGBColor(0xF2, 0x6B, 0x1D)
ACCENT_GREEN = RGBColor(0x18, 0x7A, 0x4A)
ACCENT_RED = RGBColor(0xC6, 0x3D, 0x2F)
TEXT_DARK = RGBColor(0x20, 0x25, 0x2B)
TEXT_MUTED = RGBColor(0x62, 0x6B, 0x73)
BG_LIGHT = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def parse_args():
    parser = argparse.ArgumentParser(
        description="Generate markdown and PPT artifacts for fixed-nprobe recall/QPS comparison."
    )
    parser.add_argument(
        "--summary-json",
        default="/home/v-mochengli/test/tenant0_fixed_nprobe_compare_20260417_summary.json",
        help="Path to the fixed-nprobe summary JSON.",
    )
    parser.add_argument(
        "--report-md",
        default="/home/v-mochengli/test/tenant0_fixed_nprobe_compare_20260417_report.md",
        help="Output markdown report path.",
    )
    parser.add_argument(
        "--report-ppt",
        default="/home/v-mochengli/test/tenant0_fixed_nprobe_compare_20260417_report.pptx",
        help="Output PPT path.",
    )
    return parser.parse_args()


def load_summary(path: Path):
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def group_rows(summary):
    grouped = defaultdict(list)
    for row in summary["raw_rows"]:
        grouped[(row["variant"], row["level_name"])].append(row)
    for rows in grouped.values():
        rows.sort(key=lambda item: item["nprobe"])
    return grouped


def compute_selected_rows(summary, grouped_rows):
    threshold = summary["recall_threshold"]
    selected = {}
    for variant in VARIANT_ORDER:
        for level in LEVEL_ORDER:
            rows = grouped_rows[(variant, level)]
            pick = None
            for row in rows:
                if row["recall"] >= threshold:
                    pick = dict(row)
                    pick["status"] = "selected"
                    break
            if pick is None:
                pick = dict(rows[-1])
                pick["status"] = "not-reached"
            selected[(variant, level)] = pick
    return selected


def compute_level_comparisons(selected_rows):
    comparisons = []
    for level in LEVEL_ORDER:
        old_row = selected_rows[("old_single", level)]
        current_row = selected_rows[("v5_headbundle", level)]
        qps_ratio = current_row["qps"] / old_row["qps"] if old_row["qps"] else None
        recall_gap = current_row["recall"] - old_row["recall"]
        qps_gap = current_row["qps"] - old_row["qps"]
        comparisons.append(
            {
                "level": level,
                "old": old_row,
                "current": current_row,
                "qps_ratio": qps_ratio,
                "recall_gap": recall_gap,
                "qps_gap": qps_gap,
            }
        )
    return comparisons


def fmt_float(value, digits=2):
    return f"{value:.{digits}f}"


def fmt_ratio(value):
    return "NA" if value is None else f"{value:.2f}x"


def build_markdown(summary_path: Path, report_ppt: Path, summary, comparisons):
    lead_points = []
    for item in comparisons:
        level_label = LEVEL_LABELS[item["level"]]
        ratio = item["qps_ratio"]
        if ratio is None:
            verdict = "no comparable QPS"
        elif ratio >= 1.0:
            verdict = f"current is faster ({fmt_ratio(ratio)})"
        else:
            verdict = f"current is slower ({fmt_ratio(ratio)})"
        lead_points.append(
            f"- {level_label}: {verdict}; current recall {item['current']['recall']:.4f}, original recall {item['old']['recall']:.4f}."
        )

    selected_rows = [
        "| Level | Original nprobe | Original Recall | Original QPS | Current nprobe | Current Recall | Current QPS | Current/Original QPS |",
        "| --- | ---: | ---: | ---: | ---: | ---: | ---: | ---: |",
    ]
    for item in comparisons:
        selected_rows.append(
            "| {level} | {old_nprobe} | {old_recall:.4f} | {old_qps:.2f} | {cur_nprobe} | {cur_recall:.4f} | {cur_qps:.2f} | {ratio} |".format(
                level=LEVEL_LABELS[item["level"]],
                old_nprobe=item["old"]["nprobe"],
                old_recall=item["old"]["recall"],
                old_qps=item["old"]["qps"],
                cur_nprobe=item["current"]["nprobe"],
                cur_recall=item["current"]["recall"],
                cur_qps=item["current"]["qps"],
                ratio=fmt_ratio(item["qps_ratio"]),
            )
        )

    detail_rows = [
        "| Level | Original Avg Latency ms | Current Avg Latency ms | Original Avg Posting Read | Current Avg Posting Read | QPS Delta | Recall Delta |",
        "| --- | ---: | ---: | ---: | ---: | ---: | ---: |",
    ]
    for item in comparisons:
        detail_rows.append(
            "| {level} | {old_lat:.2f} | {cur_lat:.2f} | {old_post:.2f} | {cur_post:.2f} | {qps_gap:+.2f} | {recall_gap:+.4f} |".format(
                level=LEVEL_LABELS[item["level"]],
                old_lat=item["old"]["avg_latency_ms"],
                cur_lat=item["current"]["avg_latency_ms"],
                old_post=item["old"]["avg_posting_read"],
                cur_post=item["current"]["avg_posting_read"],
                qps_gap=item["qps_gap"],
                recall_gap=item["recall_gap"],
            )
        )

    lines = [
        "# Tenant0 Fixed-Nprobe Recall/QPS Report",
        "",
        "## Experiment Setup",
        f"- Generated from: {summary_path}",
        f"- Created at: {summary['created_at']}",
        f"- Queries: {summary['num_queries']}",
        f"- TopK: {summary['topk']}",
        f"- Recall threshold: {summary['recall_threshold']}",
        "- Fixed nprobe sweep: " + ", ".join(str(item) for item in summary["nprobes"]),
        "- Baseline variant: old_single (tags_1m)",
        "- Current variant: v5_headbundle (hierarchical_tags)",
        "- Caveat: this is a level-aligned comparison, not an exact same-tag A/B.",
        "",
        "## Executive Summary",
        *lead_points,
        "",
        "## Best Operating Point Under Recall >= 95%",
        *selected_rows,
        "",
        "## Latency And Posting Cost At The Selected Point",
        *detail_rows,
        "",
        "## Interpretation",
        "- The current head-bundle path is not uniformly better: it loses on broad tags at the top of the hierarchy and wins clearly on narrower tags.",
        "- The largest regression is Org, where the current path needs nprobe 384 to cross the 95% recall bar while the original path crosses it at nprobe 64.",
        "- The largest gain is Project, where the original path needs nprobe 768 to cross the bar, while the current path crosses it at nprobe 384.",
        "- Team is the cleanest win: slightly higher recall with materially lower nprobe and 2.77x QPS.",
        "",
        "## Output Artifacts",
        f"- Raw machine summary: {summary_path}",
        f"- Presentation: {report_ppt}",
        "",
    ]
    return "\n".join(lines)


def write_markdown(path: Path, content: str):
    path.write_text(content, encoding="utf-8")


def add_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else TEXT_DARK
    subtitle_color = RGBColor(0xC8, 0xD6, 0xE5) if dark else TEXT_MUTED
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    if subtitle:
        p2 = frame.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = subtitle_color


def add_bullets(slide, items, left, top, width, height, font_size=17, color=TEXT_DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(7)


def add_table(slide, headers, rows, left, top, width, height):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(11)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER

    for row_index, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            cell = table.cell(row_index, col)
            cell.text = str(value)
            if row_index % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = TEXT_DARK
            paragraph.alignment = PP_ALIGN.CENTER


def add_clustered_bar_chart(slide, title, categories, series_specs, left, top, width, height, y_axis_title):
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top - 0.35), Inches(width), Inches(0.3))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(16)
    title_p.font.bold = True
    title_p.font.color.rgb = TEXT_DARK

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_specs:
        chart_data.add_series(name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = y_axis_title
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    for idx, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = ACCENT_BLUE if idx == 0 else ACCENT_ORANGE


def add_ratio_chart(slide, title, categories, values, left, top, width, height):
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top - 0.35), Inches(width), Inches(0.3))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(16)
    title_p.font.bold = True
    title_p.font.color.rgb = TEXT_DARK

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Current / Original QPS", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "QPS ratio"
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    fill = chart.series[0].format.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_GREEN


def add_scatter_chart(slide, title, old_rows, current_rows, left, top, width, height):
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top - 0.35), Inches(width), Inches(0.3))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(16)
    title_p.font.bold = True
    title_p.font.color.rgb = TEXT_DARK

    chart_data = XyChartData()
    old_series = chart_data.add_series(VARIANT_LABELS["old_single"])
    for row in old_rows:
        old_series.add_data_point(row["recall"], row["qps"])

    current_series = chart_data.add_series(VARIANT_LABELS["v5_headbundle"])
    for row in current_rows:
        current_series.add_data_point(row["recall"], row["qps"])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER_LINES,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Recall@10"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "QPS"
    chart.category_axis.minimum_scale = 0.30
    chart.category_axis.maximum_scale = 1.00
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)

    for idx, series in enumerate(chart.series):
        color = ACCENT_BLUE if idx == 0 else ACCENT_ORANGE
        series.format.line.color.rgb = color
        series.marker.style = XL_MARKER_STYLE.CIRCLE if idx == 0 else XL_MARKER_STYLE.DIAMOND
        series.marker.size = 7
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.marker.format.line.color.rgb = color


def build_presentation(summary_path: Path, report_md: Path, summary, grouped_rows, comparisons, output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, TITLE_BG)
    add_title(
        slide,
        "Tenant0 Fixed-Nprobe Recall/QPS Comparison",
        "Fair comparison under recall >= 95%",
        dark=True,
    )
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.7),
        Inches(11.7),
        Inches(4.6),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(0x14, 0x33, 0x57)
    hero.line.color.rgb = RGBColor(0x2C, 0x5D, 0x96)
    add_bullets(
        slide,
        [
            f"Created at: {summary['created_at']}",
            f"Queries: {summary['num_queries']} | TopK: {summary['topk']} | Recall threshold: {summary['recall_threshold']}",
            "Variants: Original single-index vs current head-bundle",
            "Note: this is level-aligned, not exact same-tag A/B.",
        ],
        left=1.2,
        top=2.2,
        width=10.8,
        height=1.6,
        font_size=18,
        color=WHITE,
    )
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(12.0), Inches(0.4))
    footer_p = footer.text_frame.paragraphs[0]
    footer_p.text = f"Summary JSON: {summary_path.name} | Markdown: {report_md.name}"
    footer_p.font.size = Pt(11)
    footer_p.font.color.rgb = RGBColor(0xC8, 0xD6, 0xE5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_title(slide, "Key Findings")
    add_bullets(
        slide,
        [
            "Org regresses materially: original crosses 95% recall at nprobe 64, current needs 384 and only delivers 0.28x QPS.",
            "Dept is close but still slower: 0.85x QPS at nearly identical recall.",
            "Team is the cleanest gain: current reaches slightly higher recall with lower nprobe and 2.77x QPS.",
            "Project is the biggest lower-level win: current crosses 95% recall at nprobe 384 while original needs 768, giving 3.57x QPS.",
            "The current scheme wins on narrow tags and still loses on broad hierarchy levels.",
        ],
        left=0.9,
        top=1.4,
        width=11.5,
        height=4.3,
    )
    methodology_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(5.5),
        Inches(11.5),
        Inches(1.1),
    )
    methodology_box.fill.solid()
    methodology_box.fill.fore_color.rgb = BG_LIGHT
    methodology_box.line.color.rgb = RGBColor(0xD9, 0xDF, 0xE7)
    method_text = slide.shapes.add_textbox(Inches(1.15), Inches(5.8), Inches(11.0), Inches(0.5))
    method_p = method_text.text_frame.paragraphs[0]
    method_p.text = "Method: fixed nprobe sweep over 64..1024, select the highest QPS point satisfying recall >= 95%."
    method_p.font.size = Pt(15)
    method_p.font.color.rgb = TEXT_DARK

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_title(slide, "Selected Operating Points")
    headers = [
        "Level",
        "Orig nprobe",
        "Orig recall",
        "Orig QPS",
        "Cur nprobe",
        "Cur recall",
        "Cur QPS",
        "QPS x",
    ]
    rows = []
    for item in comparisons:
        rows.append(
            [
                LEVEL_LABELS[item["level"]],
                item["old"]["nprobe"],
                f"{item['old']['recall']:.4f}",
                f"{item['old']['qps']:.2f}",
                item["current"]["nprobe"],
                f"{item['current']['recall']:.4f}",
                f"{item['current']['qps']:.2f}",
                fmt_ratio(item["qps_ratio"]),
            ]
        )
    add_table(slide, headers, rows, left=0.55, top=1.4, width=12.2, height=2.6)
    add_bullets(
        slide,
        [
            "Original Project only becomes comparable after pushing nprobe to 768.",
            "Current Project reaches the threshold at nprobe 384, which is the main reason for its 3.57x QPS advantage.",
            "Org remains the main gap for the head-bundle design under a strict fairness constraint.",
        ],
        left=0.8,
        top=4.6,
        width=11.7,
        height=1.8,
        font_size=15,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_title(slide, "Bar Comparison At The Selected Point")
    categories = [LEVEL_LABELS[level] for level in LEVEL_ORDER]
    add_clustered_bar_chart(
        slide,
        "QPS under recall >= 95%",
        categories,
        [
            (VARIANT_LABELS["old_single"], [selected["old"] for selected in [{"old": item['old']['qps']} for item in comparisons]]),
            (VARIANT_LABELS["v5_headbundle"], [selected["cur"] for selected in [{"cur": item['current']['qps']} for item in comparisons]]),
        ],
        left=0.6,
        top=1.5,
        width=6.0,
        height=4.7,
        y_axis_title="QPS",
    )
    add_clustered_bar_chart(
        slide,
        "Recall at the selected point",
        categories,
        [
            (VARIANT_LABELS["old_single"], [selected["old"] for selected in [{"old": item['old']['recall']} for item in comparisons]]),
            (VARIANT_LABELS["v5_headbundle"], [selected["cur"] for selected in [{"cur": item['current']['recall']} for item in comparisons]]),
        ],
        left=6.8,
        top=1.5,
        width=5.9,
        height=4.7,
        y_axis_title="Recall@10",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_title(slide, "QPS Ratio By Level")
    add_ratio_chart(
        slide,
        "Current / Original QPS at the 95% recall operating point",
        categories,
        [item["qps_ratio"] for item in comparisons],
        left=0.9,
        top=1.5,
        width=11.3,
        height=4.8,
    )
    reference = slide.shapes.add_textbox(Inches(0.95), Inches(6.45), Inches(11.0), Inches(0.35))
    ref_p = reference.text_frame.paragraphs[0]
    ref_p.text = "Reference line is 1.0x conceptually: Team and Project are above it, Org and Dept are below it."
    ref_p.font.size = Pt(13)
    ref_p.font.color.rgb = TEXT_MUTED

    for pair in [("org", "dept"), ("team", "project")]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_bg(slide, WHITE)
        title = "Recall-QPS Curves: Org and Dept" if pair[0] == "org" else "Recall-QPS Curves: Team and Project"
        add_title(slide, title)
        for index, level in enumerate(pair):
            left = 0.7 + index * 6.2
            add_scatter_chart(
                slide,
                LEVEL_LABELS[level],
                grouped_rows[("old_single", level)],
                grouped_rows[("v5_headbundle", level)],
                left=left,
                top=1.5,
                width=5.6,
                height=4.9,
            )

    prs.save(output_path)


def main():
    args = parse_args()
    summary_path = Path(args.summary_json)
    report_md = Path(args.report_md)
    report_ppt = Path(args.report_ppt)

    summary = load_summary(summary_path)
    grouped_rows = group_rows(summary)
    selected_rows = compute_selected_rows(summary, grouped_rows)
    comparisons = compute_level_comparisons(selected_rows)

    markdown = build_markdown(summary_path, report_ppt, summary, comparisons)
    write_markdown(report_md, markdown)
    build_presentation(summary_path, report_md, summary, grouped_rows, comparisons, report_ppt)

    print(f"Wrote markdown report: {report_md}")
    print(f"Wrote presentation: {report_ppt}")


if __name__ == "__main__":
    main()